# -*- coding: utf-8 -*-

from __future__ import annotations

import io
import re
from typing import Any

from django.contrib.auth import get_user_model
from django.core.files.base import ContentFile

from projects.models import Project, ProjectDocument


_SUPPORTED_KINDS = {
    "workbook",
    "run_sheet",
    "checklist",
    "slides_outline",
    "lesson_plan",
}

_DOCX_OPTIONAL_KINDS = {"workbook", "lesson_plan"}
_XLSX_OPTIONAL_KINDS = {"run_sheet", "checklist", "workbook"}
_PPTX_OPTIONAL_KINDS = {"slides_outline"}


def execute_export_capabilities() -> dict[str, bool]:
    caps = {"docx": True, "xlsx": True, "pptx": True}
    try:
        import docx  # type: ignore  # noqa: F401
    except Exception:
        caps["docx"] = False
    try:
        import openpyxl  # type: ignore  # noqa: F401
    except Exception:
        caps["xlsx"] = False
    try:
        import pptx  # type: ignore  # noqa: F401
    except Exception:
        caps["pptx"] = False
    return caps


def _safe_stem(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", str(value or "").strip()).strip("._-")
    return cleaned or "artefact"


def _slugify(value: str) -> str:
    text = re.sub(r"[^A-Za-z0-9]+", "-", str(value or "").strip().lower()).strip("-")
    return text or "untitled"


def _list_str(value: Any) -> list[str]:
    if not isinstance(value, list):
        return []
    return [str(v).strip() for v in value if str(v).strip()]


def _guess_kind_from_text(text: str) -> str:
    raw = str(text or "").strip().lower()
    if "lesson" in raw:
        return "lesson_plan"
    if "slide" in raw:
        return "slides_outline"
    if "run sheet" in raw or "run-sheet" in raw:
        return "run_sheet"
    if "checklist" in raw:
        return "checklist"
    if "workbook" in raw:
        return "workbook"
    return "workbook"


def _normalise_proposed_rows(rows: list[Any]) -> tuple[list[dict], list[str]]:
    out: list[dict] = []
    warnings: list[str] = []
    for idx, row in enumerate(rows):
        if isinstance(row, dict):
            kind = str(row.get("kind") or "").strip().lower()
            title = str(row.get("title") or "").strip()
            notes = str(row.get("notes") or "").strip()
            if not kind and title:
                kind = _guess_kind_from_text(title)
            if not title:
                title = kind.replace("_", " ").title() if kind else f"Proposed artefact {idx + 1}"
            out.append({"kind": kind, "title": title, "notes": notes})
            continue
        text_row = str(row or "").strip()
        if not text_row:
            continue
        kind = _guess_kind_from_text(text_row)
        out.append({"kind": kind, "title": text_row, "notes": ""})
        warnings.append(f"Converted string artefact proposal to object at index {idx}.")
    return out, warnings


def _resolve_uploaded_by(*, project: Project, user_id: int | None = None):
    if user_id:
        user_model = get_user_model()
        user = user_model.objects.filter(id=int(user_id)).first()
        if user is not None:
            return user
    return project.owner


def _get_project(project_id: int) -> Project:
    project = Project.objects.filter(id=int(project_id)).first()
    if project is None:
        raise ValueError("Project not found.")
    return project


def _persist_execute_artefact(
    *,
    project: Project,
    chat_id: int,
    turn_id: str,
    kind: str,
    title: str,
    ext: str,
    body: bytes,
    content_type: str,
    user_id: int | None = None,
) -> ProjectDocument:
    turn_stem = _safe_stem(turn_id)[:80]
    slug = _slugify(title)
    rel_name = f"derax/{int(chat_id)}/{turn_stem}__{kind}__{slug}.{ext}"
    uploaded_by = _resolve_uploaded_by(project=project, user_id=user_id)
    doc = ProjectDocument(
        project=project,
        title=f"DERAX EXECUTE {kind} {title}"[:200],
        original_name=rel_name[:255],
        content_type=content_type[:120],
        size_bytes=len(body),
        uploaded_by=uploaded_by,
    )
    doc.file.save(rel_name, ContentFile(body), save=False)
    doc.save()
    return doc


def _section(lines: list[str], heading: str, values: list[str]) -> None:
    lines.append(f"## {heading}")
    lines.append("")
    if values:
        for v in values:
            lines.append(f"- {v}")
    else:
        lines.append("- [TBD]")
    lines.append("")


def build_markdown_for_kind(kind: str, payload: dict, title: str, notes: str) -> str:
    kind_key = str(kind or "").strip().lower()
    intent = dict(payload.get("intent") or {})
    destination = str(intent.get("destination") or "").strip() or "[TBD]"
    success = _list_str(intent.get("success_criteria"))
    constraints = _list_str(intent.get("constraints"))
    non_goals = _list_str(intent.get("non_goals"))
    questions = _list_str(intent.get("open_questions"))

    lines = [
        f"# {str(title or '').strip() or '[TBD]'}",
        "",
        "## Destination",
        "",
        destination,
        "",
    ]
    _section(lines, "Success criteria", success)
    _section(lines, "Constraints", constraints)
    _section(lines, "Non-goals", non_goals)
    _section(lines, "Open questions", questions)

    if notes.strip():
        lines.extend(["## Notes", "", notes.strip(), ""])

    if kind_key == "workbook":
        lines.extend(
            [
                "## Session purpose",
                "",
                "- [TBD]",
                "",
                "## Agenda blocks",
                "",
                "- [TBD] 00:00-00:00",
                "",
                "## Prompt blocks by success criteria",
                "",
            ]
        )
        if success:
            for item in success:
                lines.append(f"- Prompt: How will we prove '{item}'?")
        else:
            lines.append("- Prompt: [TBD]")
        lines.extend(
            [
                "",
                "## Decisions to capture",
                "",
                "- [TBD]",
                "",
                "## Action capture",
                "",
                "| Owner | Next step | Due |",
                "| --- | --- | --- |",
                "| [TBD] | [TBD] | [TBD] |",
                "",
            ]
        )
    elif kind_key == "run_sheet":
        lines.extend(
            [
                "## Pre-reads",
                "",
                "- [TBD]",
                "",
                "## Materials checklist",
                "",
                "- [TBD]",
                "",
                "## Minute-by-minute outline",
                "",
                "- 00:00-00:00 [TBD]",
                "",
                "## Facilitator prompts",
                "",
                "- [TBD]",
                "",
                "## Decision points",
                "",
                "- [TBD]",
                "",
            ]
        )
    elif kind_key == "checklist":
        lines.extend(
            [
                "## Before",
                "",
            ]
        )
        for item in success or ["[TBD]"]:
            lines.append(f"- Confirm setup for: {item}")
        lines.extend(["", "## During", ""])
        for item in success or ["[TBD]"]:
            lines.append(f"- Capture evidence for: {item}")
        lines.extend(["", "## After", ""])
        for item in success or ["[TBD]"]:
            lines.append(f"- Verify completion for: {item}")
        lines.append("")
    elif kind_key == "slides_outline":
        lines.extend(
            [
                "## Slide outline",
                "",
                "- Slide 1: Destination overview",
                "  - " + destination,
                "- Slide 2: Success criteria",
            ]
        )
        if success:
            for item in success[:5]:
                lines.append(f"  - {item}")
        else:
            lines.append("  - [TBD]")
        lines.extend(
            [
                "- Slide 3: Constraints and non-goals",
                "  - Key constraints",
                "  - Key non-goals",
                "- Slide 4: Open questions",
                "  - Priorities to confirm",
                "",
            ]
        )
    elif kind_key == "lesson_plan":
        lines.extend(
            [
                "## Learning objectives",
                "",
            ]
        )
        for item in success or ["[TBD]"]:
            lines.append(f"- {item}")
        lines.extend(
            [
                "",
                "## Materials",
                "",
                "- [TBD]",
                "",
                "## Lesson outline",
                "",
                "- Intro: [TBD]",
                "- Activity: [TBD]",
                "- Check: [TBD]",
                "- Wrap: [TBD]",
                "",
                "## Assessment prompts",
                "",
                "- [TBD]",
                "",
            ]
        )
    else:
        lines.extend(["## Content", "", "- [TBD]", ""])

    return "\n".join(lines).strip() + "\n"


def _parse_markdown_line(line: str) -> tuple[str, str]:
    text = str(line or "").rstrip()
    if text.startswith("# "):
        return "h1", text[2:].strip()
    if text.startswith("## "):
        return "h2", text[3:].strip()
    if text.startswith("- "):
        return "li", text[2:].strip()
    if text.startswith("|"):
        return "table", text
    return "p", text.strip()


def build_docx_for_markdown(md_text: str) -> bytes:
    try:
        from docx import Document  # type: ignore
    except Exception as exc:
        raise RuntimeError(str(exc))

    document = Document()
    for raw in str(md_text or "").splitlines():
        kind, value = _parse_markdown_line(raw)
        if not value:
            continue
        if kind == "h1":
            document.add_heading(value, level=1)
        elif kind == "h2":
            document.add_heading(value, level=2)
        elif kind == "li":
            document.add_paragraph(value, style="List Bullet")
        else:
            document.add_paragraph(value)

    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def build_xlsx_for_kind(kind: str, payload: dict, title: str, notes: str) -> bytes:
    try:
        from openpyxl import Workbook  # type: ignore
    except Exception as exc:
        raise RuntimeError(str(exc))

    wb = Workbook()
    ws = wb.active
    ws.title = "Plan"
    ws.append(["Field", "Value"])

    intent = dict(payload.get("intent") or {})
    destination = str(intent.get("destination") or "").strip() or "[TBD]"
    ws.append(["Title", str(title or "").strip() or "[TBD]"])
    ws.append(["Kind", str(kind or "").strip() or "[TBD]"])
    ws.append(["Destination", destination])
    if str(notes or "").strip():
        ws.append(["Notes", str(notes).strip()])

    success = _list_str(intent.get("success_criteria"))
    constraints = _list_str(intent.get("constraints"))
    non_goals = _list_str(intent.get("non_goals"))
    open_questions = _list_str(intent.get("open_questions"))

    ws2 = wb.create_sheet("Checklist")
    ws2.append(["Group", "Item", "Status"])
    for item in success or ["[TBD]"]:
        ws2.append(["Success criteria", item, "TODO"])
    for item in constraints or ["[TBD]"]:
        ws2.append(["Constraints", item, "TODO"])
    for item in non_goals or ["[TBD]"]:
        ws2.append(["Non-goals", item, "TODO"])
    for item in open_questions or ["[TBD]"]:
        ws2.append(["Open questions", item, "TODO"])

    output = io.BytesIO()
    wb.save(output)
    return output.getvalue()


def build_pptx_for_kind(kind: str, payload: dict, title: str, notes: str) -> bytes:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        raise RuntimeError(str(exc))

    intent = dict(payload.get("intent") or {})
    destination = str(intent.get("destination") or "").strip() or "[TBD]"
    success = _list_str(intent.get("success_criteria")) or ["[TBD]"]
    constraints = _list_str(intent.get("constraints")) or ["[TBD]"]
    non_goals = _list_str(intent.get("non_goals")) or ["[TBD]"]
    open_questions = _list_str(intent.get("open_questions")) or ["[TBD]"]

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = str(title or "").strip() or "EXECUTE Output"
    slide.placeholders[1].text = destination

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Success criteria"
    slide.placeholders[1].text = "\n".join([f"- {item}" for item in success[:8]])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Constraints and non-goals"
    slide.placeholders[1].text = (
        "Constraints:\n"
        + "\n".join([f"- {item}" for item in constraints[:5]])
        + "\n\nNon-goals:\n"
        + "\n".join([f"- {item}" for item in non_goals[:5]])
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Open questions"
    slide.placeholders[1].text = "\n".join([f"- {item}" for item in open_questions[:8]])

    if str(notes or "").strip():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Notes"
        slide.placeholders[1].text = str(notes).strip()

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def generate_artefacts_from_execute_payload(
    *,
    project_id: int,
    chat_id: int,
    turn_id: str,
    payload: dict,
    user_id: int | None = None,
) -> dict:
    phase = str((payload.get("meta") or {}).get("phase") or "").strip().upper()
    if phase != "EXECUTE":
        raise ValueError("EXECUTE artefact generation requires meta.phase == EXECUTE.")

    project = _get_project(project_id)
    artefacts = dict(payload.get("artefacts") or {})
    proposed, normalise_warnings = _normalise_proposed_rows(list(artefacts.get("proposed") or []))
    generated_rows: list[dict] = []

    results_generated: list[dict] = []
    warnings: list[str] = list(normalise_warnings)
    caps = execute_export_capabilities()
    docx_available = bool(caps.get("docx"))
    xlsx_available = bool(caps.get("xlsx"))
    pptx_available = bool(caps.get("pptx"))

    for row in proposed:
        if not isinstance(row, dict):
            continue
        kind = str(row.get("kind") or "").strip().lower()
        title = str(row.get("title") or "").strip() or kind.replace("_", " ").title() or "Untitled"
        notes = str(row.get("notes") or "").strip()
        if kind not in _SUPPORTED_KINDS:
            warnings.append(f"Unsupported artefact kind: {kind or '(blank)'}")
            continue

        md_text = build_markdown_for_kind(kind, payload, title, notes)
        md_doc = _persist_execute_artefact(
            project=project,
            chat_id=chat_id,
            turn_id=turn_id,
            kind=kind,
            title=title,
            ext="md",
            body=md_text.encode("utf-8"),
            content_type="text/markdown",
            user_id=user_id,
        )
        md_entry = {"artefact_id": str(md_doc.id), "kind": kind, "title": title}
        generated_rows.append(md_entry)
        results_generated.append(md_entry)

        if kind in _DOCX_OPTIONAL_KINDS:
            if not docx_available:
                warnings.append(f"python-docx not installed; skipped docx for {kind}")
                continue
            try:
                docx_bytes = build_docx_for_markdown(md_text)
            except Exception:
                warnings.append(f"docx generation failed for {kind}; markdown created")
                continue
            docx_doc = _persist_execute_artefact(
                project=project,
                chat_id=chat_id,
                turn_id=turn_id,
                kind=kind,
                title=title,
                ext="docx",
                body=docx_bytes,
                content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                user_id=user_id,
            )
            docx_entry = {"artefact_id": str(docx_doc.id), "kind": kind, "title": f"{title} (docx)"}
            generated_rows.append(docx_entry)
            results_generated.append(docx_entry)

        if kind in _XLSX_OPTIONAL_KINDS:
            if not xlsx_available:
                warnings.append(f"openpyxl not installed; skipped xlsx for {kind}")
            else:
                try:
                    xlsx_bytes = build_xlsx_for_kind(kind, payload, title, notes)
                except Exception:
                    warnings.append(f"xlsx generation failed for {kind}; markdown created")
                else:
                    xlsx_doc = _persist_execute_artefact(
                        project=project,
                        chat_id=chat_id,
                        turn_id=turn_id,
                        kind=kind,
                        title=title,
                        ext="xlsx",
                        body=xlsx_bytes,
                        content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        user_id=user_id,
                    )
                    xlsx_entry = {"artefact_id": str(xlsx_doc.id), "kind": kind, "title": f"{title} (xlsx)"}
                    generated_rows.append(xlsx_entry)
                    results_generated.append(xlsx_entry)

        if kind in _PPTX_OPTIONAL_KINDS:
            if not pptx_available:
                warnings.append(f"python-pptx not installed; skipped pptx for {kind}")
            else:
                try:
                    pptx_bytes = build_pptx_for_kind(kind, payload, title, notes)
                except Exception:
                    warnings.append(f"pptx generation failed for {kind}; markdown created")
                else:
                    pptx_doc = _persist_execute_artefact(
                        project=project,
                        chat_id=chat_id,
                        turn_id=turn_id,
                        kind=kind,
                        title=title,
                        ext="pptx",
                        body=pptx_bytes,
                        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        user_id=user_id,
                    )
                    pptx_entry = {"artefact_id": str(pptx_doc.id), "kind": kind, "title": f"{title} (pptx)"}
                    generated_rows.append(pptx_entry)
                    results_generated.append(pptx_entry)

    artefacts["proposed"] = proposed
    artefacts["generated"] = generated_rows
    payload["artefacts"] = artefacts
    return {"generated": results_generated, "warnings": warnings}
